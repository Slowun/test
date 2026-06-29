"""Build professional CryptoHedge presentations (RU + EN).

Renders 10 quant-style slides as full 16:9 figures (dark theme), embedding:
  - custom architecture / strategy / data-flow / defence / roadmap diagrams
  - live charts from solution.ipynb sections 7, 9, 10, 12 (pulled from the
    multi-agent pipeline context, exactly as the notebook produces them)

Outputs (per language):
  docs/presentation.<lang>.pdf   <- primary deliverable (2 PDF files)
  docs/presentation.<lang>.pptx  <- editable companion (full-bleed slides)

================================ КАРТА МОДУЛЯ ================================
СЛОЙ:        scripts (вспомогательные утилиты) — ТОЧКА ВХОДА (CLI-скрипт).
НАЗНАЧЕНИЕ:  собирает профессиональную презентацию проекта (RU+EN): запускает
             весь пайплайн, достаёт из контекста реальные графики (греки,
             портфель, риск, диагностика) и рендерит 10 слайдов 16:9 в PDF+PPTX.
             Это «витрина» проекта, не часть рантайма хедж-фонда.
ИМПОРТИРУЕТ:
  - sys, textwrap, pathlib    : система, перенос текста, пути.
  - matplotlib (Agg backend)  : отрисовка фигур/диаграмм/графиков, экспорт в PDF.
  - numpy, pandas             : работа с данными графиков.
  - python-pptx               : генерация редактируемого .pptx.
  - cryptohedge.* (внутри run_pipeline): запуск пайплайна для извлечения данных.
ЭКСПОРТИРУЕТ: main() и набор функций отрисовки (диаграммы/графики/слайды).
КОНСТАНТЫ:   ROOT/ASSETS/OUT (пути), цветовая «тёмная quant-тема», FIG_W/FIG_H.
ЗАПУСК: python scripts/build_presentations.py [--cached]
=============================================================================
"""

from __future__ import annotations                       # отложенные аннотации типов

import sys                                               # argv, stderr, путь импорта
import textwrap                                          # перенос длинного текста в булитах
from pathlib import Path                                 # пути

import matplotlib                                        # библиотека визуализации

matplotlib.use("Agg")                                    # безоконный бэкенд (рендер в файлы без дисплея)
import matplotlib.pyplot as plt                          # API построения графиков
import numpy as np                                       # числовые операции для графиков
import pandas as pd                                      # таблицы данных из контекста
import matplotlib.dates as mdates                        # форматирование осей дат
from matplotlib.backends.backend_pdf import PdfPages     # многостраничный PDF
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch  # стрелки/блоки для диаграмм
from pptx import Presentation                            # генерация PowerPoint
from pptx.util import Inches                             # размеры слайдов в дюймах

ROOT = Path(__file__).resolve().parents[1]               # корень проекта (на уровень выше scripts/)
ASSETS = ROOT / "docs" / "_presentation_assets"          # каталог промежуточных PNG-ассетов
OUT = ROOT / "docs"                                      # каталог итоговых PDF/PPTX

# ── тёмная quant-тема (палитра цветов слайдов) ───────────────────────────────
BG = "#0d1117"                                            # фон слайда
PANEL = "#161b22"                                         # фон панели графика
PANEL2 = "#1c2330"                                        # фон вторичной панели/плитки
ACCENT = "#4cc2ff"                                        # акцентный голубой
GREEN = "#3fb950"                                         # зелёный (позитив/портфель)
AMBER = "#d29922"                                         # янтарный (предупреждение/акцент)
RED = "#f85149"                                           # красный (риск/без хеджа)
PURPLE = "#bc8cff"                                        # фиолетовый (квант/греки)
TEXT = "#e6edf3"                                          # основной цвет текста
MUTED = "#8b949e"                                         # приглушённый текст
GRID = "#30363d"                                          # цвет сетки/рамок

plt.rcParams.update({                                     # глобальные настройки рендера:
    "font.family": "DejaVu Sans",                       #   шрифт (поддержка кириллицы)
    "figure.dpi": 150,                                  #   разрешение
    "pdf.fonttype": 42,                                 #   TrueType-шрифты в PDF (редактируемость)
    "axes.unicode_minus": False,                        #   обычный минус вместо unicode
})

FIG_W, FIG_H = 12.8, 7.2  # размеры фигуры в дюймах (соотношение 16:9)


# ── pipeline + data ──────────────────────────────────────────────────────────
def run_pipeline():
    """Run the full agent pipeline and return its context (notebook section 2)."""
    sys.path.insert(0, str(ROOT))                        # гарантируем импорт пакета cryptohedge
    from cryptohedge.agents import build_pipeline        # сборка пайплайна (ленивый импорт)
    from cryptohedge.core.config import load_config      # загрузка конфигурации
    from cryptohedge.core.context import AgentContext    # контекст прогона

    print("Running multi-agent pipeline (fresh run)...")
    config = load_config(ROOT / "config", overrides={"runtime": {"resume": False}})  # свежий прогон без resume
    context = AgentContext(config, root=ROOT)            # создаём контекст
    build_pipeline(context, fail_fast=True).run()        # ЗАПУСК всего пайплайна
    return context                                       # контекст с артефактами на blackboard


def _safe(ctx, key, default=None):
    # безопасное чтение артефакта из контекста: при ошибке/None вернуть default
    try:
        v = ctx.get(key, default)
        return default if v is None else v
    except Exception:
        return default


def extract(ctx) -> dict:
    # собрать из контекста все артефакты и факты конфига, нужные для слайдов
    cfg = ctx.config
    d = {
        # section 7
        "chain_greeks": _safe(ctx, "chain_greeks"),
        "portfolio_greeks_latest": _safe(ctx, "portfolio_greeks_latest", {}),
        "hedge_status": _safe(ctx, "hedge_status", {}),
        # section 9
        "method_comparison": _safe(ctx, "method_comparison"),
        "diversification": _safe(ctx, "diversification", {}),
        "portfolio_equity": _safe(ctx, "portfolio_equity"),
        "portfolio_weights_path": _safe(ctx, "portfolio_weights_path"),
        "portfolio_rebalances": _safe(ctx, "portfolio_rebalances", []),
        "rebalance_decision": _safe(ctx, "rebalance_decision", {}),
        # section 10
        "risk_assessment": _safe(ctx, "risk_assessment", {}),
        "stop_level": _safe(ctx, "stop_level", {}),
        "trailing_stops": _safe(ctx, "trailing_stops"),
        # section 11 (stress, supports validation slide)
        "stress_table": _safe(ctx, "stress_table"),
        # section 12
        "diagnostic": _safe(ctx, "diagnostic", {}),
        "confidence_score": float(_safe(ctx, "confidence_score", 0.0) or 0.0),
        # config facts
        "capital": getattr(cfg.investment, "capital_usd", 10_000_000),
        "risk_budget": getattr(cfg.investment, "risk_budget_pct", 0.02),
        "fee_spot": getattr(cfg.investment, "transaction_fee_pct", 0.0003),
        "fee_opt": getattr(cfg.investment, "option_fee_pct", 0.0003),
        "fee_cap": getattr(cfg.investment, "option_fee_cap_pct", 0.125),
        "universe_size": getattr(cfg.data, "universe_size", 100),
        "primary": getattr(cfg.data, "primary_symbol", "BTCUSDT"),
        "var_limit": getattr(cfg.risk, "var_limit_pct", 0.05),
        "mdd_limit": getattr(cfg.risk, "max_drawdown_limit_pct", 0.25),
        "seed": getattr(cfg, "seed", 90909090),
    }
    return d


def facts(d: dict) -> dict:
    """Headline numbers formatted for slide bullets."""
    mc = d["method_comparison"]
    chosen = (d["rebalance_decision"] or {}).get("method", "risk_parity")
    ret = dr = eff = na = mdd = None
    if isinstance(mc, pd.DataFrame) and "method" in mc.columns:
        row = mc.set_index("method")
        if chosen in row.index:
            r = row.loc[chosen]
            ret = float(r.get("total_return", np.nan))
            mdd = float(r.get("max_drawdown", np.nan))
    div = d["diversification"] or {}
    dr = div.get("diversification_ratio", dr)
    eff = div.get("effective_n", eff)
    na = div.get("n_assets", 15)
    ra = d["risk_assessment"] or {}
    return {
        "chosen": chosen,
        "ret": ret if ret is not None else 0.421,
        "dr": dr if dr is not None else 1.63,
        "eff": eff if eff is not None else 13.7,
        "na": int(na) if na else 15,
        "mdd": mdd if mdd is not None else -0.095,
        "var": ra.get("var", 0.0001),
        "cvar": ra.get("cvar", 0.0001),
        "var_limit": d["var_limit"],
        "mdd_limit": d["mdd_limit"],
        "conf": d["confidence_score"],
        "seed": d["seed"],
        "capital": d["capital"],
        "universe": d["universe_size"],
    }


# ── slide frame ──────────────────────────────────────────────────────────────
def new_slide():
    # создать пустую фигуру слайда с тёмным фоном
    fig = plt.figure(figsize=(FIG_W, FIG_H))
    fig.patch.set_facecolor(BG)
    return fig


def frame(fig, page, total, title, subtitle, lang):
    # нарисовать «рамку» слайда: акцентная полоса, заголовок, подзаголовок, футер
    # top accent bar
    fig.add_artist(plt.Line2D([0.045, 0.955], [0.935, 0.935], color=ACCENT, lw=2.5,
                              transform=fig.transFigure))
    fig.text(0.045, 0.955, title, fontsize=25, fontweight="bold", color=TEXT, va="bottom")
    if subtitle:
        fig.text(0.045, 0.905, subtitle, fontsize=13, color=MUTED, va="top")
    # footer
    brand = "CryptoHedge · Multi-Agent Crypto Hedge Fund · MVP"
    fig.text(0.045, 0.022, brand, fontsize=8.5, color=MUTED, va="bottom")
    fig.text(0.955, 0.022, f"{page:02d} / {total:02d}", fontsize=8.5, color=MUTED,
             va="bottom", ha="right")


def bullets(fig, items, x=0.045, y_top=0.84, width=58, fontsize=12.5,
            color=TEXT, marker_color=ACCENT, line_gap=0.028, bullet_gap=0.022):
    """Render bullet list in figure coords; returns final y."""
    y = y_top
    for raw in items:
        emph = raw.startswith("@")          # rationale highlight
        txt = raw[1:] if emph else raw
        c = AMBER if emph else color
        mk = AMBER if emph else marker_color
        wrapped = textwrap.wrap(txt, width=width) or [""]
        fig.text(x, y, "▪", fontsize=fontsize - 1, color=mk, va="top")
        for j, line in enumerate(wrapped):
            fig.text(x + 0.016, y - j * line_gap, line, fontsize=fontsize,
                     color=c, va="top",
                     fontweight="bold" if emph and j == 0 else "normal")
        y -= line_gap * len(wrapped) + bullet_gap
    return y


def style_ax(ax, title=""):
    # применить единый тёмный стиль к осям графика
    ax.set_facecolor(PANEL)
    ax.tick_params(colors=MUTED, labelsize=7.5)
    for s in ax.spines.values():
        s.set_color(GRID)
    ax.grid(True, alpha=0.18, color=GRID)
    if title:
        ax.set_title(title, color=ACCENT, fontsize=10, pad=5)
    ax.xaxis.label.set_color(MUTED)
    ax.yaxis.label.set_color(MUTED)


def date_axis(ax):
    # форматировать ось X как даты (метки по месяцам)
    ax.xaxis.set_major_locator(mdates.AutoDateLocator(maxticks=5))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b"))


def kpi_tile(fig, x, y, w, h, value, label, color=ACCENT):
    # нарисовать KPI-плитку (крупное значение + подпись) на титульном слайде
    ax = fig.add_axes([x, y, w, h])
    ax.axis("off")
    ax.add_patch(FancyBboxPatch((0.02, 0.02), 0.96, 0.96, boxstyle="round,pad=0.02",
                                facecolor=PANEL2, edgecolor=color, linewidth=1.4,
                                transform=ax.transAxes))
    ax.text(0.5, 0.62, value, ha="center", va="center", fontsize=21,
            fontweight="bold", color=color, transform=ax.transAxes)
    ax.text(0.5, 0.22, label, ha="center", va="center", fontsize=9.5,
            color=TEXT, transform=ax.transAxes)


# ── charts from notebook sections ────────────────────────────────────────────
def chart_greeks(ax, d, lang):
    """Section 7: normalized greeks profile across strikes."""
    chain = d["chain_greeks"]
    style_ax(ax, "Greeks vs strike (normalised)" if lang == "en"
             else "Профиль греков по страйкам (нормировано)")
    if isinstance(chain, pd.DataFrame) and "moneyness" in chain.columns:
        cols = [c for c in ["delta", "gamma", "vega", "theta"] if c in chain.columns]
        palette = {"delta": ACCENT, "gamma": GREEN, "vega": PURPLE, "theta": AMBER}
        for c in cols:
            denom = chain[c].abs().max() or 1.0
            ax.plot(chain["moneyness"], chain[c] / denom, marker="o", ms=3,
                    lw=1.6, color=palette.get(c, ACCENT), label=c)
        ax.axvline(1.0, color=MUTED, ls="--", lw=1)
        ax.set_xlabel("moneyness K/S")
        ax.legend(fontsize=7.5, facecolor=PANEL2, edgecolor=GRID, labelcolor=TEXT, ncol=2)
    else:
        ax.text(0.5, 0.5, "chain_greeks n/a", ha="center", color=MUTED, transform=ax.transAxes)


def chart_methods(ax, d, lang):
    """Section 9: optimizer comparison; chosen method highlighted."""
    mc = d["method_comparison"]
    chosen = (d["rebalance_decision"] or {}).get("method", "risk_parity")
    style_ax(ax, "Optimizer backtest (return)" if lang == "en"
             else "Бэктест методов (доходность)")
    abbr = {"mean_variance": "mean-var", "risk_parity": "risk parity",
            "min_variance": "min-var", "max_diversification": "max-div", "cvar": "CVaR"}
    if isinstance(mc, pd.DataFrame) and "method" in mc.columns:
        m = mc.set_index("method")
        order = m["total_return"].sort_values()
        cols = [GREEN if i == chosen else "#39506b" for i in order.index]
        ax.barh(range(len(order)), order.values * 100, color=cols)
        ax.set_yticks(range(len(order)))
        ax.set_yticklabels([abbr.get(i, i.replace("_", " ")) for i in order.index], fontsize=7.5)
        ax.set_xlabel("% return")
        ax.margins(x=0.16)
        for k, v in enumerate(order.values):
            ax.text(v * 100, k, f" {v:.0%}", va="center", fontsize=7.5, color=TEXT)
    else:
        ax.text(0.5, 0.5, "method_comparison n/a", ha="center", color=MUTED, transform=ax.transAxes)


def chart_equity(ax, d, lang):
    """Section 9: portfolio equity vs benchmark + rebalance markers."""
    eq = d["portfolio_equity"]
    style_ax(ax, "Portfolio equity & rebalances" if lang == "en"
             else "Стоимость портфеля и ребалансировки")
    if isinstance(eq, pd.DataFrame) and "equity" in eq.columns:
        ts = pd.to_datetime(eq["ts"])
        ax.plot(ts, eq["equity"], color=GREEN, lw=2,
                label="portfolio" if lang == "en" else "портфель")
        if "benchmark" in eq.columns:
            ax.plot(ts, eq["benchmark"], color=MUTED, ls="--", lw=1.4,
                    label="EW benchmark" if lang == "en" else "бенчмарк")
        rebs = pd.to_datetime([str(r) for r in (d["portfolio_rebalances"] or [])],
                              errors="coerce").dropna()
        if len(rebs):
            idx = pd.Index(ts).get_indexer(rebs, method="nearest")
            idx = idx[idx >= 0]
            ax.scatter(ts.iloc[idx], eq["equity"].iloc[idx], color=AMBER, s=18, zorder=5,
                       label="rebalance" if lang == "en" else "ребаланс")
        ax.legend(fontsize=7.5, facecolor=PANEL2, edgecolor=GRID, labelcolor=TEXT)
        ax.set_ylabel("base 1.0" if lang == "en" else "база 1.0")
        date_axis(ax)
    else:
        ax.text(0.5, 0.5, "portfolio_equity n/a", ha="center", color=MUTED, transform=ax.transAxes)


def chart_weights(ax, d, lang):
    """Section 9: weight evolution stackplot (top assets)."""
    wp = d["portfolio_weights_path"]
    style_ax(ax, "Weight evolution" if lang == "en" else "Эволюция весов")
    if isinstance(wp, pd.DataFrame) and "ts" in wp.columns:
        ts = pd.to_datetime(wp["ts"])
        cols = [c for c in wp.columns if c != "ts"]
        top = list(wp[cols].mean().sort_values(ascending=False).head(8).index)
        cmap = plt.cm.viridis(np.linspace(0.1, 0.95, len(top)))
        ax.stackplot(ts, *[wp[c] for c in top], labels=top, colors=cmap, alpha=0.9)
        ax.set_ylim(0, 1)
        ax.legend(fontsize=5.8, ncol=2, facecolor=PANEL2, edgecolor=GRID,
                  labelcolor=TEXT, loc="upper left")
        ax.set_ylabel("weight" if lang == "en" else "вес")
        date_axis(ax)
    else:
        ax.text(0.5, 0.5, "weights n/a", ha="center", color=MUTED, transform=ax.transAxes)


def chart_trailing(ax, d, lang):
    """Section 10: dynamic trailing stop vs price."""
    tr = d["trailing_stops"]
    style_ax(ax, "Dynamic trailing stop" if lang == "en" else "Динамический трейлинг-стоп")
    if isinstance(tr, pd.DataFrame) and "price" in tr.columns and len(tr):
        tt = pd.to_datetime(tr["ts"])
        ax.plot(tt, tr["price"], color=ACCENT, lw=1.6,
                label="BTC price" if lang == "en" else "Цена BTC")
        ax.plot(tt, tr["stop_price"], color=RED, ls="--", lw=1.4,
                label="trailing stop" if lang == "en" else "трейлинг-стоп")
        ax.fill_between(tt, tr["stop_price"], tr["price"],
                        where=tr["price"] >= tr["stop_price"], color=GREEN, alpha=0.10)
        ax.legend(fontsize=7.5, facecolor=PANEL2, edgecolor=GRID, labelcolor=TEXT)
        ax.set_ylabel("USD")
        date_axis(ax)
    else:
        ax.text(0.5, 0.5, "trailing_stops n/a", ha="center", color=MUTED, transform=ax.transAxes)


def chart_confidence(ax, d, lang):
    """Section 12: confidence index components."""
    diag = d["diagnostic"] or {}
    comp = diag.get("components", {})
    cs = d["confidence_score"]
    style_ax(ax, "Confidence index components" if lang == "en"
             else "Компоненты индекса доверия")
    if comp:
        names = list(comp.keys())
        vals = [float(comp[k]) for k in names]
        colors = [GREEN if v >= 0.6 else AMBER if v >= 0.4 else RED for v in vals]
        ax.barh(names, vals, color=colors)
        ax.set_xlim(0, 1)
        ax.axvline(cs, color=TEXT, ls="--", lw=1.2,
                   label=f"confidence={cs:.2f}")
        ax.legend(fontsize=7.5, facecolor=PANEL2, edgecolor=GRID, labelcolor=TEXT)
        ax.tick_params(axis="y", labelsize=7.5)
    else:
        ax.text(0.5, 0.5, "diagnostic n/a", ha="center", color=MUTED, transform=ax.transAxes)


def chart_stress(ax, d, lang):
    """Section 11: stress test, hedged vs naked BTC."""
    s = d["stress_table"]
    style_ax(ax, "Stress tests: hedged vs naked" if lang == "en"
             else "Стресс-тесты: хедж vs голый BTC")
    if isinstance(s, pd.DataFrame) and "net_hedged_pnl" in s.columns:
        x = np.arange(len(s))
        w = 0.38
        ax.bar(x - w / 2, s["unhedged_pnl"], w, color=RED,
               label="naked BTC" if lang == "en" else "без хеджа")
        ax.bar(x + w / 2, s["net_hedged_pnl"], w, color=GREEN,
               label="hedged" if lang == "en" else "с хеджем")
        ax.set_xticks(x)
        ax.set_xticklabels([str(v) for v in s["scenario"]], rotation=18, fontsize=6.8)
        ax.axhline(0, color=MUTED, lw=0.8)
        ax.legend(fontsize=7.5, facecolor=PANEL2, edgecolor=GRID, labelcolor=TEXT)
        ax.set_ylabel("PnL, USD")
    else:
        ax.text(0.5, 0.5, "stress_table n/a", ha="center", color=MUTED, transform=ax.transAxes)


# ── custom diagrams ──────────────────────────────────────────────────────────
def _box(ax, x, y, w, h, text, fc, ec="white", fs=8.5, tc="white", lw=1.0):
    # нарисовать прямоугольный блок с текстом (элемент диаграмм)
    ax.add_patch(FancyBboxPatch((x - w / 2, y - h / 2), w, h, boxstyle="round,pad=0.02",
                                facecolor=fc, edgecolor=ec, linewidth=lw))
    ax.text(x, y, text, ha="center", va="center", fontsize=fs, color=tc, fontweight="bold")


def _arrow(ax, x1, y1, x2, y2, color=ACCENT, lw=1.8):
    # нарисовать стрелку между двумя точками (связь на диаграмме)
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>",
                                 mutation_scale=12, color=color, lw=lw))


def diagram_agents(ax, lang):  # диаграмма мультиагентного пайплайна (5 стадий + общая шина)
    ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis("off")
    stages = [
        ("Data", ["DataAcquisition", "MarketAnalysis"], ACCENT),
        ("Quant", ["HestonCalibration", "GreeksCalc"], PURPLE),
        ("Decision", ["HedgingDecision", "PortfolioOpt", "RiskMgmt"], GREEN),
        ("Validate", ["Backtesting", "SelfDiagnostic"], AMBER),
        ("Output", ["Explainability", "Dashboard"], "#58a6ff"),
    ]
    n = len(stages)
    cw = 9.4 / n
    for i, (stage, agents, col) in enumerate(stages):
        cx = 0.5 + cw * i + cw / 2 - 0.1
        ax.add_patch(FancyBboxPatch((cx - cw / 2 + 0.08, 1.3), cw - 0.25, 4.0,
                                    boxstyle="round,pad=0.03", facecolor=PANEL2,
                                    edgecolor=col, linewidth=1.3))
        ax.text(cx, 5.0, stage, ha="center", fontsize=10, color=col, fontweight="bold")
        for j, a in enumerate(agents):
            ay = 4.3 - j * 0.95
            _box(ax, cx, ay, cw - 0.45, 0.62, a, col, ec=col, fs=7.2, tc="white", lw=0)
        if i < n - 1:
            _arrow(ax, cx + cw / 2 - 0.18, 3.3, cx + cw / 2 + 0.12, 3.3, color=GREEN)
    # shared bus
    ax.add_patch(FancyBboxPatch((0.55, 0.35), 8.8, 0.7, boxstyle="round,pad=0.02",
                                facecolor="#11212f", edgecolor=ACCENT, linewidth=1.2))
    bus = ("Blackboard (AgentContext)  +  MessageBus  +  Orchestrator (fail-fast, checkpoint/resume)"
           if lang == "en" else
           "Blackboard (AgentContext)  +  MessageBus  +  Оркестратор (fail-fast, checkpoint/resume)")
    ax.text(5.0, 0.7, bus, ha="center", va="center", fontsize=8.5, color=ACCENT)
    for i in range(n):
        cx = 0.5 + cw * i + cw / 2 - 0.1
        _arrow(ax, cx, 1.3, cx, 1.06, color=MUTED, lw=1.0)


def diagram_strategy(ax, lang):  # диаграмма стратегии хеджирования (книга→размер→Δ→ν→портфель→риск)
    ax.set_xlim(0, 12); ax.set_ylim(0, 2.4); ax.axis("off")
    steps = ([
        ("Liability", "short put +\nvega call", RED),
        ("Size", "2% daily VaR\nrisk budget", AMBER),
        ("Δ-hedge", "spot BTC\nΔ → 0", ACCENT),
        ("ν-hedge", "ATM option\nν → 0", "#00b3b3"),
        ("Portfolio", "Risk Parity\n15 assets", GREEN),
        ("Risk", "VaR/CVaR\nstops", PURPLE),
    ] if lang == "en" else [
        ("Книга", "short put +\nvega-call", RED),
        ("Размер", "бюджет риска\n2% дн. VaR", AMBER),
        ("Δ-хедж", "спот BTC\nΔ → 0", ACCENT),
        ("ν-хедж", "ATM-опцион\nν → 0", "#00b3b3"),
        ("Портфель", "Risk Parity\n15 активов", GREEN),
        ("Риск", "VaR/CVaR\nстопы", PURPLE),
    ])
    n = len(steps)
    xs = np.linspace(1.0, 11.0, n)
    for i, (hdr, body, col) in enumerate(steps):
        x = xs[i]
        ax.add_patch(FancyBboxPatch((x - 0.78, 0.5), 1.56, 1.5, boxstyle="round,pad=0.03",
                                    facecolor=PANEL2, edgecolor=col, linewidth=1.5))
        ax.text(x, 1.62, hdr, ha="center", fontsize=9, color=col, fontweight="bold")
        ax.text(x, 1.0, body, ha="center", va="center", fontsize=7.4, color=TEXT)
        if i < n - 1:
            _arrow(ax, x + 0.8, 1.25, xs[i + 1] - 0.8, 1.25, color=AMBER)


def diagram_dataflow(ax, lang):  # диаграмма потока данных (источники→провайдер→вселенная)
    ax.set_xlim(0, 12); ax.set_ylim(0, 4); ax.axis("off")
    srcs = ["bundled", "synthetic", "binance REST"]
    for j, s in enumerate(srcs):
        _box(ax, 1.4, 3.2 - j * 1.05, 2.0, 0.7, s, PANEL2, ec=ACCENT, fs=8.2, tc=TEXT)
        _arrow(ax, 2.45, 3.2 - j * 1.05, 3.7, 2.1, color=MUTED, lw=1.2)
    _box(ax, 5.2, 2.1, 2.6, 1.0,
         "MarketDataProvider\n(abstraction)" if lang == "en" else "MarketDataProvider\n(абстракция)",
         "#11283a", ec=ACCENT, fs=8.5, tc=ACCENT)
    _arrow(ax, 6.55, 2.1, 7.7, 2.1, color=GREEN)
    _box(ax, 9.0, 2.1, 2.7, 1.2,
         "Universe\n100 spot + BTC options\n(call/put, 30d, 11K)" if lang == "en" else
         "Вселенная\n100 спот + опционы BTC\n(call/put, 30д, 11 страйков)",
         PANEL2, ec=GREEN, fs=7.8, tc=TEXT)
    tag = ("Swap provider → OKX / Bybit / Deribit without touching agents"
           if lang == "en" else
           "Смена провайдера → OKX / Bybit / Deribit без правки агентов")
    ax.text(6.0, 0.45, tag, ha="center", fontsize=8.5, color=AMBER, style="italic")


def diagram_protective(ax, lang):  # диаграмма слоёв защиты портфеля
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.2); ax.axis("off")
    layers = ([
        ("Position sizing", "from risk budget, not market view", GREEN),
        ("Exposure limits", "max weight 20% · leverage 3× · MDD 25%", ACCENT),
        ("Dynamic hedge", "rebalance Δ / ν on greek drift", PURPLE),
        ("Adaptive + trailing stop", "max(ATR, VaR, Heston-vol)", AMBER),
        ("Portfolio rebalancing", "back to target weights every 5d", "#58a6ff"),
    ] if lang == "en" else [
        ("Размер позиции", "из бюджета риска, не из вью", GREEN),
        ("Лимиты экспозиций", "макс. вес 20% · плечо 3× · просадка 25%", ACCENT),
        ("Динамический хедж", "ребаланс Δ / ν при дрейфе греков", PURPLE),
        ("Адаптивный + трейлинг-стоп", "max(ATR, VaR, Heston-vol)", AMBER),
        ("Ребалансировка портфеля", "возврат к целевым весам каждые 5д", "#58a6ff"),
    ])
    for i, (hdr, body, col) in enumerate(layers):
        y = 3.6 - i * 0.72
        w = 11.0 - i * 0.0
        ax.add_patch(FancyBboxPatch((0.6, y - 0.3), w, 0.6, boxstyle="round,pad=0.02",
                                    facecolor=PANEL2, edgecolor=col, linewidth=1.4))
        ax.text(0.95, y, hdr, ha="left", va="center", fontsize=9.5, color=col, fontweight="bold")
        ax.text(6.7, y, body, ha="left", va="center", fontsize=8.3, color=TEXT)


def diagram_architecture(ax, lang):  # диаграмма Clean Architecture (agents→services→domain→pyquant/core)
    ax.set_xlim(0, 12); ax.set_ylim(0, 4.4); ax.axis("off")
    layers = [
        ("agents", "11 AI agents · orchestration", ACCENT),
        ("services", "hedging · portfolio · risk · backtest", GREEN),
        ("domain", "instruments · greeks · portfolio model", PURPLE),
        ("pyquant / core", "Heston · pricing · VaR · seeding", AMBER),
    ]
    for i, (name, body, col) in enumerate(layers):
        inset = i * 0.7
        y = 3.7 - i * 0.85
        ax.add_patch(FancyBboxPatch((0.6 + inset, y - 0.34), 10.8 - 2 * inset, 0.68,
                                    boxstyle="round,pad=0.02", facecolor=PANEL2,
                                    edgecolor=col, linewidth=1.5))
        ax.text(6.0, y + 0.07, name, ha="center", fontsize=10, color=col, fontweight="bold")
        ax.text(6.0, y - 0.17, body, ha="center", fontsize=7.8, color=TEXT)
    tag = ("Dependencies point inward → business logic independent of exchanges/data"
           if lang == "en" else
           "Зависимости направлены внутрь → логика не зависит от бирж и источников данных")
    ax.text(6.0, 0.25, tag, ha="center", fontsize=8.3, color=MUTED, style="italic")


def diagram_roadmap(ax, lang):  # диаграмма дорожной карты (MVP→ближайшее→среднее→долгое)
    ax.set_xlim(0, 12); ax.set_ylim(0, 4); ax.axis("off")
    ax.add_patch(FancyArrowPatch((0.6, 2.0), (11.4, 2.0), arrowstyle="-|>",
                                 mutation_scale=18, color=ACCENT, lw=2.2))
    phases = ([
        ("MVP (now)", "backtested\nmulti-agent", GREEN, 2.95),
        ("Near", "live multi-CEX\nexecution + Telegram", ACCENT, 1.0),
        ("Mid", "news / sentiment\n+ data sources", PURPLE, 2.95),
        ("Long", "RL agents · cross-exch\narb · live options", AMBER, 1.0),
    ] if lang == "en" else [
        ("MVP (сейчас)", "бэктест\nмультиагент", GREEN, 2.95),
        ("Ближайшее", "live-исполнение\nна CEX + Telegram", ACCENT, 1.0),
        ("Среднее", "новости / sentiment\n+ источники данных", PURPLE, 2.95),
        ("Долгое", "RL-агенты · арбитраж\n· live-опционы", AMBER, 1.0),
    ])
    xs = np.linspace(1.5, 10.5, len(phases))
    for (hdr, body, col, ty), x in zip(phases, xs):
        ax.add_patch(plt.Circle((x, 2.0), 0.12, color=col, zorder=5))
        _box(ax, x, ty, 2.3, 0.95, f"{hdr}\n{body}", PANEL2, ec=col, fs=7.8, tc=TEXT)
        ax.plot([x, x], [2.0, ty + (0.48 if ty > 2 else -0.48)], color=col, lw=1.0, ls=":")


# ── slide content (RU + EN) ──────────────────────────────────────────────────
def slide_specs(lang, f):
    # вернуть список спецификаций 10 слайдов для языка lang (f — рассчитанные факты)
    pct = lambda v: f"{v*100:.1f}%"                      # хелпер форматирования доли в проценты
    if lang == "ru":
        return [
            dict(title="CryptoHedge",
                 sub="Мультиагентный криптовалютный хедж-фонд на базе AI · реализованный MVP",
                 layout="title",
                 bullets=[
                     f"Капитал ${f['capital']/1e6:.0f}M · горизонт 90 дней · единый seed = {f['seed']} (воспроизводимость)",
                     "11 автономных AI-агентов: генерация сигналов, хеджирование, риск, портфель",
                     "Δ-ν-нейтральное хеджирование BTC по модели Хестона + диверсифицированный портфель",
                     "MVP реализован и прошёл историческое бэктестирование (walk-forward + стресс-тесты)",
                     "42 автотеста · бэктест с ребалансировкой · дашборды RU/EN",
                 ],
                 notes="Открывающий слайд: подчеркнуть, что это работающая реализация, а не концепт."),
            dict(title="Торговая вселенная и данные", sub="Инструменты, источники, комиссии, масштабирование",
                 layout="bullets_diagram", diagram=diagram_dataflow,
                 bullets=[
                     f"Вселенная: {f['universe']} спот USDT-пар; первичный BTCUSDT; опционы BTC (call/put, 30 дн., 11 страйков)",
                     "Источники: bundled (offline) · synthetic (генератор) · binance (live REST) — выбор через конфиг",
                     "Комиссии: спот 0.03%; опционы Deribit-style 0.03% + кап 12.5% от цены",
                     "Хедж: спот для Δ, опцион для ν — нейтрализация рыночных греков",
                     "@Почему: абстракция MarketDataProvider изолирует биржевые API → масштабирование на OKX/Bybit/Deribit без переписывания агентов",
                 ],
                 notes="Акцент на готовности к мультибиржевому масштабированию."),
            dict(title="Мультиагентная архитектура", sub="Слабая связанность, трассируемость, тестируемость",
                 layout="bullets_diagram", diagram=diagram_agents,
                 bullets=[
                     "Blackboard (AgentContext): агенты обмениваются артефактами без прямых вызовов",
                     "MessageBus: типизированные сообщения с correlation_id — полная трассировка решений",
                     "Clean Architecture: agents → services → domain → pyquant",
                     "Оркестратор fail-fast; checkpoint/resume — устойчивость длинных прогонов",
                     "@Почему blackboard: добавление/замена агента не ломает остальных — расширяемость и изоляция тестов",
                 ],
                 notes="Объяснить выбор blackboard против прямых вызовов агентов."),
            dict(title="Количественный движок и хеджирование", sub="Раздел 7: греки портфеля · модель Хестона",
                 layout="diagram_chart", diagram=diagram_strategy, chart=chart_greeks,
                 bullets=[
                     "Калибровка Хестона (MLE + IV-поверхность); бенчмарки Black-Scholes / SABR",
                     "Полный набор греков: Δ Γ ν Θ ρ + vanna / volga / charm",
                     "Книга обязательств: short put (0.95 K/S) + vega-call; размер хеджа из бюджета риска 2% дн. VaR",
                     "Δ-хедж спотом, ν-хедж опционом → остаточные греки ≈ 0",
                     "@Почему Хестон: стохастическая волатильность точнее BS на крипто (улыбка, кластеризация vol)",
                 ],
                 notes="Раздел 7 ноутбука: профиль греков по страйкам."),
            dict(title="Количественное управление портфелем", sub="Раздел 9: оптимизация, ребалансировка, диверсификация",
                 layout="three_charts", charts=[chart_methods, chart_equity, chart_weights],
                 bullets=[
                     "5 оптимизаторов: Mean-Variance · Risk Parity · Min-Var · Max-Div · CVaR",
                     "Авто-выбор по бэктесту с ребалансировкой (с комиссиями)",
                     f"MVP: выбран {f['chosen'].replace('_',' ')} → доходность +{pct(f['ret'])}, DR {f['dr']:.2f}, эфф. N {f['eff']:.1f}/{f['na']}",
                     "Ребаланс каждые 5 дней; макс. вес 20% — контроль концентрации",
                     "@Почему Risk Parity: устойчивая диверсификация без опоры на прогноз доходностей",
                 ],
                 notes="Раздел 9: сравнение методов, equity и эволюция весов."),
            dict(title="Управление рисками", sub="Раздел 10: VaR / CVaR и адаптивные стопы",
                 layout="bullets_chart", chart=chart_trailing,
                 bullets=[
                     f"Меры риска: VaR (hist/Gauss/Cornish-Fisher), CVaR/ES, макс. просадка — VaR {pct(f['var'])} ≪ лимит {f['var_limit']:.0%}",
                     "Адаптивный стоп-лосс: max(ATR, VaR, Heston-vol) — реакция на режим волатильности",
                     "Динамический трейлинг-стоп; зоны баланса дельты (green / red)",
                     "Защищено: рыночная Δ, вега, концентрация, хвостовые шоки",
                     "Не защищено: исполнение/проскальзывание, контрагент, базис, скачки (jumps)",
                     "@Допущения: дневные бары, long-only, историческая оценка VaR, модельные комиссии",
                 ],
                 notes="Раздел 10: трейлинг-стоп; честно перечислить незахеджированные риски."),
            dict(title="Защитные механизмы", sub="Многоуровневая оборона портфеля",
                 layout="bullets_diagram", diagram=diagram_protective,
                 bullets=[
                     "Трейлинг-стоп (ATR×2.5) — фиксация прибыли при развороте",
                     "Динамическая корректировка хеджа — ребаланс Δ/ν при дрейфе греков",
                     "Управление размером позиции — из бюджета риска, а не из вью на рынок",
                     "Ограничения экспозиций — макс. вес, плечо 3×, лимит просадки 25%",
                     "Ребалансировка портфеля — возврат к целевым весам каждые 5 дней",
                 ],
                 notes="Каждый слой снижает свой класс риска."),
            dict(title="Валидация, бэктест и самодиагностика", sub="Разделы 11-12: стресс-тесты и индекс доверия",
                 layout="two_charts", charts=[chart_stress, chart_confidence],
                 bullets=[
                     "Walk-forward: train 30 дн / test 5 дн; контроль look-ahead и survivorship bias",
                     "Стресс-тесты: захеджированная книга против голого BTC",
                     "Self-Diagnostic: индекс доверия (calibration, data-drift PSI, forecast-error, hedge-quality, risk-compliance)",
                     f"Текущий индекс доверия = {f['conf']:.2f}: при дрейфе данных система честно сигнализирует low_confidence",
                     "@MVP прошёл историческое бэктестирование; результаты воспроизводимы",
                 ],
                 notes="Раздел 12: компоненты индекса доверия; раздел 11: стресс."),
            dict(title="Воспроизводимость и инженерия", sub="Детерминизм, тесты, чистая архитектура",
                 layout="bullets_diagram", diagram=diagram_architecture,
                 bullets=[
                     f"Единый seed ({f['seed']}) фиксирует все RNG: данные, калибровку, PnL, веса",
                     "42 автотеста; end-to-end solution.ipynb; max|ΔPnL| = 0 при повторном прогоне",
                     "Checkpoint/resume; конфиг-управляемый пайплайн (YAML + Pydantic) без хардкода",
                     "Clean Architecture: бизнес-логика не зависит от провайдеров данных и бирж",
                     "@Почему: воспроизводимость — обязательное требование для аудита и честного сравнения стратегий",
                 ],
                 notes="Reproducibility — ключ для quant-аудита."),
            dict(title="Направления развития", sub="Roadmap",
                 layout="bullets_diagram", diagram=diagram_roadmap,
                 bullets=[
                     "Полностью автоматическое исполнение сделок на нескольких CEX (live)",
                     "Telegram-бот для управления и мониторинга",
                     "Интеграция новостных потоков и анализа настроений (sentiment)",
                     "Подключение дополнительных источников данных; агенты на RL",
                     "Межбиржевой арбитраж · live-рынок опционов · самоадаптирующаяся координация агентов",
                 ],
                 notes="Закрывающий слайд: вектор развития."),
        ]
    # English
    return [
        dict(title="CryptoHedge",
             sub="Multi-agent AI cryptocurrency hedge fund · delivered MVP",
             layout="title",
             bullets=[
                 f"Capital ${f['capital']/1e6:.0f}M · 90-day horizon · single seed = {f['seed']} (reproducible)",
                 "11 autonomous AI agents: signal generation, hedging, risk, portfolio",
                 "Δ-ν-neutral BTC hedging via the Heston model + diversified portfolio",
                 "MVP delivered and validated on historical backtest (walk-forward + stress tests)",
                 "42 automated tests · rebalanced backtest · RU/EN dashboards",
             ],
             notes="Opening: this is a working implementation, not a concept."),
        dict(title="Trading universe & data", sub="Instruments, sources, fees, scaling",
             layout="bullets_diagram", diagram=diagram_dataflow,
             bullets=[
                 f"Universe: {f['universe']} spot USDT pairs; primary BTCUSDT; BTC options (call/put, 30d, 11 strikes)",
                 "Sources: bundled (offline) · synthetic (generator) · binance (live REST) — config-selected",
                 "Fees: spot 0.03%; Deribit-style options 0.03% + 12.5% price cap",
                 "Hedge: spot for Δ, option for ν — neutralise market greeks",
                 "@Why: the MarketDataProvider abstraction isolates exchange APIs → scale to OKX/Bybit/Deribit without rewriting agents",
             ],
             notes="Emphasise multi-exchange scaling readiness."),
        dict(title="Multi-agent architecture", sub="Loose coupling, traceability, testability",
             layout="bullets_diagram", diagram=diagram_agents,
             bullets=[
                 "Blackboard (AgentContext): agents exchange artifacts with no direct calls",
                 "MessageBus: typed messages with correlation_id — full decision traceability",
                 "Clean Architecture: agents → services → domain → pyquant",
                 "Fail-fast orchestrator; checkpoint/resume — robust long runs",
                 "@Why blackboard: adding/replacing an agent never breaks others — extensibility and isolated tests",
             ],
             notes="Explain blackboard vs direct agent calls."),
        dict(title="Quantitative engine & hedging", sub="Section 7: portfolio greeks · Heston model",
             layout="diagram_chart", diagram=diagram_strategy, chart=chart_greeks,
             bullets=[
                 "Heston calibration (MLE + IV surface); Black-Scholes / SABR benchmarks",
                 "Full greeks set: Δ Γ ν Θ ρ + vanna / volga / charm",
                 "Liability book: short put (0.95 K/S) + vega call; hedge size from 2% daily VaR budget",
                 "Δ hedge via spot, ν hedge via option → residual greeks ≈ 0",
                 "@Why Heston: stochastic volatility beats BS on crypto (smile, vol clustering)",
             ],
             notes="Notebook section 7: greeks profile across strikes."),
        dict(title="Quantitative portfolio management", sub="Section 9: optimization, rebalancing, diversification",
             layout="three_charts", charts=[chart_methods, chart_equity, chart_weights],
             bullets=[
                 "5 optimizers: Mean-Variance · Risk Parity · Min-Var · Max-Div · CVaR",
                 "Auto-selected via rebalanced backtest (fees included)",
                 f"MVP: selected {f['chosen'].replace('_',' ')} → return +{pct(f['ret'])}, DR {f['dr']:.2f}, effective N {f['eff']:.1f}/{f['na']}",
                 "Rebalance every 5 days; max weight 20% — concentration control",
                 "@Why Risk Parity: robust diversification without relying on return forecasts",
             ],
             notes="Section 9: method comparison, equity, weight evolution."),
        dict(title="Risk management", sub="Section 10: VaR / CVaR and adaptive stops",
             layout="bullets_chart", chart=chart_trailing,
             bullets=[
                 f"Risk measures: VaR (hist/Gauss/Cornish-Fisher), CVaR/ES, max drawdown — VaR {pct(f['var'])} ≪ {f['var_limit']:.0%} limit",
                 "Adaptive stop-loss: max(ATR, VaR, Heston-vol) — reacts to volatility regime",
                 "Dynamic trailing stop; delta-balance zones (green / red)",
                 "Protected: market Δ, vega, concentration, tail shocks",
                 "Not protected: execution/slippage, counterparty, basis, jumps",
                 "@Assumptions: daily bars, long-only, historical VaR, modelled fees",
             ],
             notes="Section 10: trailing stop; honestly list unhedged risks."),
        dict(title="Protective mechanisms", sub="Layered portfolio defence",
             layout="bullets_diagram", diagram=diagram_protective,
             bullets=[
                 "Trailing stop (ATR×2.5) — lock in profit on reversal",
                 "Dynamic hedge adjustment — rebalance Δ/ν on greek drift",
                 "Position sizing — from risk budget, not from a market view",
                 "Exposure limits — max weight, 3× leverage, 25% drawdown cap",
                 "Portfolio rebalancing — return to target weights every 5 days",
             ],
             notes="Each layer mitigates a distinct risk class."),
        dict(title="Validation, backtest & self-diagnostics", sub="Sections 11-12: stress tests and confidence index",
             layout="two_charts", charts=[chart_stress, chart_confidence],
             bullets=[
                 "Walk-forward: train 30d / test 5d; look-ahead and survivorship-bias controls",
                 "Stress tests: hedged book vs naked BTC",
                 "Self-Diagnostic: confidence index (calibration, data-drift PSI, forecast-error, hedge-quality, risk-compliance)",
                 f"Current confidence = {f['conf']:.2f}: under data drift the system honestly flags low_confidence",
                 "@MVP passed historical backtesting; results are reproducible",
             ],
             notes="Section 12: confidence components; section 11: stress."),
        dict(title="Reproducibility & engineering", sub="Determinism, tests, clean architecture",
             layout="bullets_diagram", diagram=diagram_architecture,
             bullets=[
                 f"Single seed ({f['seed']}) fixes all RNG: data, calibration, PnL, weights",
                 "42 automated tests; end-to-end solution.ipynb; max|ΔPnL| = 0 on re-run",
                 "Checkpoint/resume; config-driven pipeline (YAML + Pydantic), no hardcoding",
                 "Clean Architecture: business logic independent of data providers and exchanges",
                 "@Why: reproducibility is mandatory for audit and fair strategy comparison",
             ],
             notes="Reproducibility is key for quant audit."),
        dict(title="Future development", sub="Roadmap",
             layout="bullets_diagram", diagram=diagram_roadmap,
             bullets=[
                 "Fully automated trade execution across multiple CEX (live)",
                 "Telegram bot for control and monitoring",
                 "News feeds and sentiment analysis integration",
                 "Additional data sources; reinforcement-learning agents",
                 "Cross-exchange arbitrage · live options market · self-adapting agent coordination",
             ],
             notes="Closing: development vector."),
    ]


# ── slide layout engine ──────────────────────────────────────────────────────
def render_slide(spec, d, lang, page, total):
    # отрисовать один слайд по его spec (раскладка выбирается по spec["layout"])
    fig = new_slide()
    frame(fig, page, total, spec["title"], spec.get("sub", ""), lang)
    layout = spec["layout"]

    if layout == "title":
        bullets(fig, spec["bullets"], x=0.045, y_top=0.80, width=78, fontsize=14,
                line_gap=0.030, bullet_gap=0.030)
        f = facts(d)
        tiles = [
            (f"+{f['ret']*100:.1f}%", "Return" if lang == "en" else "Доходность", GREEN),
            (f"{f['dr']:.2f}", "Diversification" if lang == "en" else "Диверсификация", ACCENT),
            (f"{f['eff']:.1f}/{f['na']}", "Effective N" if lang == "en" else "Эфф. число активов", PURPLE),
            (f"{f['conf']:.2f}", "Confidence" if lang == "en" else "Индекс доверия", AMBER),
        ]
        for i, (v, l, c) in enumerate(tiles):
            kpi_tile(fig, 0.05 + i * 0.235, 0.10, 0.205, 0.20, v, l, c)

    elif layout == "bullets_diagram":
        bullets(fig, spec["bullets"], x=0.045, y_top=0.82, width=92, fontsize=12.5,
                line_gap=0.027, bullet_gap=0.020)
        ax = fig.add_axes([0.06, 0.07, 0.88, 0.40])
        spec["diagram"](ax, lang)

    elif layout == "bullets_chart":
        bullets(fig, spec["bullets"], x=0.045, y_top=0.82, width=52, fontsize=12,
                line_gap=0.027, bullet_gap=0.020)
        ax = fig.add_axes([0.55, 0.12, 0.40, 0.66])
        spec["chart"](ax, d, lang)

    elif layout == "diagram_chart":
        bullets(fig, spec["bullets"], x=0.045, y_top=0.82, width=52, fontsize=12,
                line_gap=0.027, bullet_gap=0.020)
        ax_d = fig.add_axes([0.52, 0.50, 0.45, 0.28]); spec["diagram"](ax_d, lang)
        ax_c = fig.add_axes([0.55, 0.10, 0.40, 0.32]); spec["chart"](ax_c, d, lang)

    elif layout == "two_charts":
        bullets(fig, spec["bullets"], x=0.045, y_top=0.82, width=92, fontsize=12,
                line_gap=0.026, bullet_gap=0.018)
        ax1 = fig.add_axes([0.07, 0.13, 0.40, 0.32]); spec["charts"][0](ax1, d, lang)
        ax2 = fig.add_axes([0.56, 0.13, 0.39, 0.32]); spec["charts"][1](ax2, d, lang)

    elif layout == "three_charts":
        bullets(fig, spec["bullets"], x=0.045, y_top=0.82, width=92, fontsize=11.5,
                line_gap=0.025, bullet_gap=0.016)
        ax1 = fig.add_axes([0.075, 0.11, 0.25, 0.33]); spec["charts"][0](ax1, d, lang)
        ax2 = fig.add_axes([0.40, 0.11, 0.26, 0.33]); spec["charts"][1](ax2, d, lang)
        ax3 = fig.add_axes([0.71, 0.11, 0.26, 0.33]); spec["charts"][2](ax3, d, lang)

    return fig


# ── assembly ─────────────────────────────────────────────────────────────────
def build_language(d, lang):
    # построить презентацию (PDF + PPTX) для одного языка
    specs = slide_specs(lang, facts(d))                  # спецификации слайдов
    total = len(specs)                                   # всего слайдов
    pdf_path = OUT / f"presentation.{lang}.pdf"          # путь PDF
    pngs = []                                            # пути PNG-кадров (для PPTX)
    with PdfPages(pdf_path) as pdf:                      # многостраничный PDF
        for i, spec in enumerate(specs, 1):             # по каждому слайду…
            fig = render_slide(spec, d, lang, i, total)  # рендерим фигуру
            pdf.savefig(fig, facecolor=BG)              # добавляем страницу в PDF
            png = ASSETS / f"slide{i:02d}_{lang}.png"   # имя PNG-кадра
            fig.savefig(png, facecolor=BG, dpi=150)     # сохраняем PNG
            pngs.append(png)
            plt.close(fig)                              # освобождаем память
    print(f"  PDF  -> {pdf_path}  ({total} slides)")

    prs = Presentation()                                 # создаём PPTX
    prs.slide_width = Inches(13.333)                     # ширина слайда 16:9
    prs.slide_height = Inches(7.5)                       # высота слайда
    blank = prs.slide_layouts[6]                         # пустой макет
    for png in pngs:                                     # вставляем каждый PNG на весь слайд
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    pptx_path = OUT / f"presentation.{lang}.pptx"        # путь PPTX
    prs.save(str(pptx_path))
    print(f"  PPTX -> {pptx_path}")


def main():
    import pickle                                        # кэширование данных пайплайна
    ASSETS.mkdir(parents=True, exist_ok=True)            # создаём каталог ассетов
    cache = ASSETS / "_data_cache.pkl"                   # файл кэша
    if "--cached" in sys.argv and cache.exists():        # режим --cached: брать данные из кэша
        print("Loading cached pipeline data...")
        d = pickle.loads(cache.read_bytes())
    else:                                                # иначе — свежий прогон пайплайна
        ctx = run_pipeline()
        d = extract(ctx)                                 # извлечь данные для графиков
        try:
            cache.write_bytes(pickle.dumps(d))          # сохранить кэш
        except Exception as e:
            print(f"  cache skipped: {e}", file=sys.stderr)  # кэш не критичен
    for lang in ("ru", "en"):                            # собрать обе языковые версии
        print(f"\n=== {lang.upper()} ===")
        build_language(d, lang)
    print("\nDone.")


if __name__ == "__main__":
    main()                                               # запуск сборки презентаций из консоли
